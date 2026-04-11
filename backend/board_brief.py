"""
Fourseat - Fourseat Decks
Auto-generates board-ready presentations from manual data input.
Produces a PowerPoint (.pptx) and a structured JSON report.
"""

import os
import json
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import anthropic
from dotenv import load_dotenv

load_dotenv()

DATA_DIR    = Path(__file__).parent.parent / "data"
OUTPUT_DIR  = DATA_DIR / "outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

anthropic_client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY", ""))

# ── Colors ────────────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x0D, 0x0D)
ACCENT     = RGBColor(0xFF, 0x6B, 0x35)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xAA, 0xAA, 0xAA)
MID_GRAY   = RGBColor(0x33, 0x33, 0x33)


# ── Slide helpers ─────────────────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, text: str, left, top, width, height,
              font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txBox


def _add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ── Slide builders ────────────────────────────────────────────────────────────

def _slide_title(prs, company: str, period: str, tagline: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, DARK_BG)
    # accent bar
    _add_rect(slide, Inches(0), Inches(3.2), Inches(10), Inches(0.06), ACCENT)
    _add_text(slide, company.upper(), Inches(0.6), Inches(1.2), Inches(8), Inches(1.2),
              font_size=44, bold=True, color=WHITE)
    _add_text(slide, f"Board Update  ·  {period}", Inches(0.6), Inches(2.5), Inches(8), Inches(0.5),
              font_size=18, color=ACCENT)
    _add_text(slide, tagline, Inches(0.6), Inches(3.5), Inches(8), Inches(1),
              font_size=14, color=LIGHT_GRAY)
    _add_text(slide, "BOARDROOM AI", Inches(7.5), Inches(6.8), Inches(2), Inches(0.4),
              font_size=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)


def _slide_metrics(prs, metrics: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DARK_BG)
    _add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.8), MID_GRAY)
    _add_text(slide, "KEY METRICS", Inches(0.4), Inches(0.15), Inches(9), Inches(0.5),
              font_size=16, bold=True, color=ACCENT)

    items = list(metrics.items())[:6]
    cols, rows = 3, 2
    for i, (k, v) in enumerate(items):
        col = i % cols
        row = i // cols
        lft = Inches(0.3 + col * 3.2)
        top = Inches(1.1 + row * 2.5)
        _add_rect(slide, lft, top, Inches(2.9), Inches(2.1), MID_GRAY)
        _add_text(slide, str(v), lft + Inches(0.15), top + Inches(0.2),
                  Inches(2.6), Inches(1.1), font_size=32, bold=True, color=ACCENT)
        _add_text(slide, k, lft + Inches(0.15), top + Inches(1.3),
                  Inches(2.6), Inches(0.6), font_size=12, color=LIGHT_GRAY)


def _slide_bullets(prs, title: str, bullets: list[str], accent_color=ACCENT):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DARK_BG)
    _add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.8), MID_GRAY)
    _add_text(slide, title.upper(), Inches(0.4), Inches(0.15), Inches(9), Inches(0.5),
              font_size=16, bold=True, color=accent_color)
    for i, b in enumerate(bullets[:6]):
        top = Inches(1.1 + i * 0.9)
        _add_rect(slide, Inches(0.4), top + Inches(0.15), Inches(0.06), Inches(0.4), accent_color)
        _add_text(slide, b, Inches(0.65), top, Inches(9), Inches(0.8), font_size=14, color=WHITE)


def _slide_narrative(prs, title: str, body: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DARK_BG)
    _add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.8), MID_GRAY)
    _add_text(slide, title.upper(), Inches(0.4), Inches(0.15), Inches(9), Inches(0.5),
              font_size=16, bold=True, color=ACCENT)
    _add_text(slide, body, Inches(0.6), Inches(1.1), Inches(8.8), Inches(5.4),
              font_size=13, color=WHITE)


# ── AI narrative generator ────────────────────────────────────────────────────

def _generate_narratives(data: dict) -> dict:
    prompt = (
        f"You are writing a board update for {data.get('company_name', 'the company')}.\n"
        f"Metrics: {json.dumps(data.get('metrics', {}))}\n"
        f"Highlights: {data.get('highlights', '')}\n"
        f"Challenges: {data.get('challenges', '')}\n"
        f"Ask of the board: {data.get('ask', '')}\n\n"
        f"Return ONLY a JSON object with these keys (no markdown):\n"
        f"{{\n"
        f'  "executive_summary": "2-3 sentence summary of the period",\n'
        f'  "wins": ["win 1", "win 2", "win 3"],\n'
        f'  "risks": ["risk 1", "risk 2", "risk 3"],\n'
        f'  "next_quarter": ["priority 1", "priority 2", "priority 3"],\n'
        f'  "ask_narrative": "one paragraph on what you need from the board"\n'
        f"}}"
    )
    try:
        resp = anthropic_client.messages.create(
            model="claude-3-haiku-20240307",
            max_tokens=1024,
            messages=[{"role": "user", "content": prompt}],
        )
        raw = resp.content[0].text.strip().lstrip("```json").lstrip("```").rstrip("```").strip()
        return json.loads(raw)
    except Exception:
        return {
            "executive_summary": data.get("highlights", "Board update for the period."),
            "wins": [data.get("highlights", "See attached")],
            "risks": [data.get("challenges", "See attached")],
            "next_quarter": ["Continue execution"],
            "ask_narrative": data.get("ask", ""),
        }


# ── Main deck builder ─────────────────────────────────────────────────────────

def generate_board_deck(data: dict) -> dict:
    """
    Generate a board deck from structured data.
    data keys: company_name, period, tagline, metrics (dict),
               highlights, challenges, ask
    Returns: dict with file path and summary
    """
    narratives = _generate_narratives(data)

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    company = data.get("company_name", "Company")
    period  = data.get("period", datetime.now().strftime("%B %Y"))

    _slide_title(prs, company, period, narratives["executive_summary"])
    _slide_metrics(prs, data.get("metrics", {}))
    _slide_bullets(prs, "Wins This Period", narratives["wins"], ACCENT)
    _slide_bullets(prs, "Risks & Challenges", narratives["risks"], RGBColor(0xFF, 0x4C, 0x4C))
    _slide_bullets(prs, "Next Quarter Priorities", narratives["next_quarter"], RGBColor(0x4C, 0xAF, 0xFF))
    _slide_narrative(prs, "Ask of the Board", narratives["ask_narrative"])

    filename  = f"{company.replace(' ', '_')}_{period.replace(' ', '_')}_BoardDeck.pptx"
    out_path  = OUTPUT_DIR / filename
    prs.save(str(out_path))

    return {
        "success": True,
        "filepath": str(out_path),
        "filename": filename,
        "narratives": narratives,
        "slides": 6,
    }
